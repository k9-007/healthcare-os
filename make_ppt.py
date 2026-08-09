"""One-shot generator for HealthcareOS.pptx (uses python-pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import os

ASSETS = "/Users/eloelo/.cursor/projects/Users-eloelo-healthcare-os/assets"
IMGS = {
    "dashboard": f"{ASSETS}/image-a14f2233-8a72-46d3-a772-21b84899993c.png",
    "patients": f"{ASSETS}/image-8f913afd-4aa0-41ed-b695-b8294e82ec3d.png",
    "careplan": f"{ASSETS}/image-d339fad9-2f63-44ff-a637-214995a873f5.png",
    "calls": f"{ASSETS}/image-6b08c9b7-551b-40ff-b9f4-a0c9f316d1ef.png",
    "caregraph": f"{ASSETS}/image-ed8cc033-0680-47cd-b631-5fcd02c33e51.png",
}

GREEN = RGBColor(0x0F, 0x9D, 0x6E)
DARK = RGBColor(0x1E, 0x29, 0x3B)
SLATE = RGBColor(0x47, 0x55, 0x69)
LIGHT = RGBColor(0xF6, 0xF9, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xDC, 0x26, 0x26)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, bold, color, space_after_pt)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (text, size, bold, color, space) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        r = p.add_run()
        r.text = text
        f = r.font
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = color
        f.name = "Avenir Next"
    return tb


def accent_bar(slide, x, y, w=Inches(0.6), h=Inches(0.07)):
    r = slide.shapes.add_shape(1, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = GREEN
    r.line.fill.background()
    r.shadow.inherit = False


def picture_fit(slide, path, x, y, max_w, max_h, border=True):
    im = Image.open(path)
    ar = im.width / im.height
    w, h = max_w, int(max_w / ar)
    if h > max_h:
        h = max_h
        w = int(max_h * ar)
    px = x + (max_w - w) // 2
    py = y + (max_h - h) // 2
    pic = slide.shapes.add_picture(path, px, py, width=w, height=h)
    if border:
        pic.line.color.rgb = RGBColor(0xD8, 0xE2, 0xDE)
        pic.line.width = Pt(1)
    return pic


# ---------- Slide 1: Title ----------
s = add_slide(DARK)
band = s.shapes.add_shape(1, 0, Inches(6.9), SW, Inches(0.6))
band.fill.solid(); band.fill.fore_color.rgb = GREEN; band.line.fill.background(); band.shadow.inherit = False
textbox(s, Inches(1.0), Inches(2.35), Inches(11.3), Inches(3.0), [
    ("HealthcareOS", 54, True, WHITE, 6),
    ("The AI Care Coordination Layer for Bharat", 26, False, RGBColor(0x9E, 0xE6, 0xC9), 18),
    ("An AI voice nurse that calls every discharged patient in their own language —\nreminds, listens, escalates, and closes the loop with the doctor.", 17, False, RGBColor(0xC7, 0xD2, 0xE0), 0),
])

# ---------- Slide 2: Problem ----------
s = add_slide()
accent_bar(s, Inches(0.9), Inches(0.75))
textbox(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(1.0), [
    ("The problem: care ends at the hospital gate", 32, True, DARK, 0),
])
probs = [
    ("Post-discharge is a black hole", "Once a patient leaves, doctors have no visibility into whether medicines are taken or symptoms worsen."),
    ("Follow-up doesn't scale", "Manual phone follow-ups are expensive; nurses can't call hundreds of patients twice a day."),
    ("Apps don't work for Bharat", "Most patients won't install an app or read English SMS — but everyone answers a phone call in their language."),
    ("Missed signals become emergencies", "Skipped doses and early symptoms (chest pain, dizziness) surface only at readmission."),
]
x0, y0, cw, ch, gap = Inches(0.9), Inches(2.1), Inches(5.6), Inches(2.1), Inches(0.35)
for i, (t, d) in enumerate(probs):
    x = x0 + (i % 2) * (cw + gap)
    y = y0 + (i // 2) * (ch + gap)
    card = s.shapes.add_shape(1, x, y, cw, ch)
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = RGBColor(0xDD, 0xE6, 0xE1); card.line.width = Pt(0.75)
    card.shadow.inherit = False
    textbox(s, x + Inches(0.3), y + Inches(0.25), cw - Inches(0.6), ch - Inches(0.5), [
        (t, 17, True, GREEN, 6),
        (d, 13.5, False, SLATE, 0),
    ])

# ---------- Slide 3: Solution ----------
s = add_slide()
accent_bar(s, Inches(0.9), Inches(0.75))
textbox(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(1.0), [
    ("The solution: a closed-loop AI voice nurse", 32, True, DARK, 0),
])
steps = [
    ("1 · Plan", "Doctor saves a care plan — medicines, dose times, follow-up questions, call window."),
    ("2 · Call", "The scheduler auto-places real phone calls at each dose time, speaking in the patient's language (Hindi, Tamil, ...)."),
    ("3 · Listen", "The patient just talks. Speech is transcribed and turned into structured data: took medicine, symptoms, urgency."),
    ("4 · Escalate", "High-urgency replies (e.g. chest pain) raise an escalation on the doctor's dashboard within seconds."),
    ("5 · Close the loop", "The doctor types advice once — it's translated and delivered back as a voice callback. No app, no SMS."),
]
y = Inches(2.05)
for t, d in steps:
    chip = s.shapes.add_shape(1, Inches(0.9), y, Inches(2.2), Inches(0.78))
    chip.fill.solid(); chip.fill.fore_color.rgb = GREEN; chip.line.fill.background(); chip.shadow.inherit = False
    tf = chip.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = t; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE
    textbox(s, Inches(3.4), y + Inches(0.02), Inches(9.0), Inches(0.9), [
        (d, 14, False, SLATE, 0),
    ], anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.99)

# ---------- Screenshot slides ----------
def shot_slide(title, sub, img, bullets):
    s = add_slide()
    accent_bar(s, Inches(0.7), Inches(0.5))
    textbox(s, Inches(0.7), Inches(0.68), Inches(12.0), Inches(0.9), [
        (title, 26, True, DARK, 2),
        (sub, 14, False, SLATE, 0),
    ])
    picture_fit(s, img, Inches(0.7), Inches(1.75), Inches(8.4), Inches(5.3))
    lines = []
    for b in bullets:
        lines.append(("•  " + b, 13, False, SLATE, 10))
    textbox(s, Inches(9.4), Inches(1.95), Inches(3.4), Inches(5.0), lines)
    return s

shot_slide(
    "Command center",
    "One glance tells the care team who needs attention right now.",
    IMGS["dashboard"],
    [
        "Live KPIs: adherence %, missed doses, patients at risk, call success.",
        "Escalation feed sorted by urgency — chest pain surfaces at the top.",
        "One-click acknowledge keeps the queue clean.",
        "14-day adherence trend across all active patients.",
    ],
)

shot_slide(
    "Patient roster",
    "Every discharged patient, their language, adherence and next scheduled call.",
    IMGS["patients"],
    [
        "Filter by Active / At-risk / Recovered.",
        "Per-patient adherence bar and risk status computed automatically.",
        "\u201cNext call\u201d shows exactly when the AI nurse dials next.",
        "Multilingual by design — Hindi, Tamil and more per patient.",
    ],
)

shot_slide(
    "Care plan → Care+ engine",
    "The doctor writes the plan once; the scheduler does the rest, every day.",
    IMGS["careplan"],
    [
        "Medicines with dose, times and instructions (\u201cBefore Food\u201d).",
        "Follow-up questions the AI nurse asks by voice.",
        "Call window guarantees no calls at odd hours (patient-local time).",
        "Edits re-sync future slots only — placed calls are never rewritten.",
    ],
)

shot_slide(
    "Calls & the closed loop",
    "Real Plivo phone calls in, structured clinical data out.",
    IMGS["calls"],
    [
        "Each call is transcribed and parsed: took medicine, symptom, urgency.",
        "\u201cChest pain, urgency High\u201d auto-raises an escalation mid-call.",
        "Doctor replies in the reply box — delivered as a voice callback in Hindi.",
        "Simulation mode lets you demo the whole flow without telephony.",
    ],
)

shot_slide(
    "Care graph — the patient's full journey",
    "A single timeline from discharge to recovery: calls, symptoms, alerts, advice.",
    IMGS["caregraph"],
    [
        "Every event is logged: care calls, symptom reports, escalations, callbacks.",
        "Doctors see cause and effect — advice sent, next call outcome.",
        "Perfect audit trail for handovers and clinical review.",
    ],
)

# ---------- Under the hood ----------
s = add_slide()
accent_bar(s, Inches(0.9), Inches(0.75))
textbox(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(1.0), [
    ("Under the hood", 32, True, DARK, 0),
])
cols = [
    ("Voice & language — Sarvam AI", [
        "LLM for call scripts & reply understanding",
        "TTS / STT in Indian languages",
        "Translate + language ID for the closed loop",
        "Deterministic fallbacks if the API is down",
    ]),
    ("Telephony — Plivo", [
        "Real outbound calls with recorded replies",
        "Answer / recording / hangup webhooks",
        "Live voice websocket streaming",
        "Simulation mode for demos & tests",
    ]),
    ("Engine — FastAPI + SQLite", [
        "Cron scheduler: materialize slots, group by time, dial",
        "Retries with backoff, call windows, missed-dose events",
        "Brain: document Q&A with citations (cite-or-refuse)",
        "React frontend, 5-language UI (en/hi/kn/mr/ta)",
    ]),
]
x = Inches(0.9)
for title, items in cols:
    card = s.shapes.add_shape(1, x, Inches(2.1), Inches(3.75), Inches(4.4))
    card.fill.solid(); card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = RGBColor(0xDD, 0xE6, 0xE1); card.line.width = Pt(0.75)
    card.shadow.inherit = False
    lines = [(title, 16, True, GREEN, 10)]
    for it in items:
        lines.append(("•  " + it, 12.5, False, SLATE, 8))
    textbox(s, x + Inches(0.28), Inches(2.35), Inches(3.2), Inches(4.0), lines)
    x += Inches(3.95)

# ---------- Closing ----------
s = add_slide(DARK)
textbox(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(2.6), [
    ("Verified end-to-end", 34, True, WHITE, 10),
    ("Plan → scheduled call → patient reply → structured data → escalation → doctor advice → voice callback.", 17, False, RGBColor(0xC7, 0xD2, 0xE0), 16),
    ("HealthcareOS — no app, no SMS. Just a phone call, in the patient's language.", 17, True, RGBColor(0x9E, 0xE6, 0xC9), 0),
])
band = s.shapes.add_shape(1, 0, Inches(6.9), SW, Inches(0.6))
band.fill.solid(); band.fill.fore_color.rgb = GREEN; band.line.fill.background(); band.shadow.inherit = False

out = "/Users/eloelo/healthcare-os/HealthcareOS.pptx"
prs.save(out)
print("saved", out, os.path.getsize(out), "bytes")
